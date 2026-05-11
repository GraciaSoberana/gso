"""
pdf_to_pptx.py
Convierte un PDF (exportado desde Advanced Slides) a PPTX.
Cada página del PDF se convierte en una imagen y se inserta como slide.

Uso:
    python pdf_to_pptx.py sermon.pdf
    python pdf_to_pptx.py sermon.pdf -o salida.pptx
    python pdf_to_pptx.py sermon.pdf --dpi 200

Requiere:
    pip install pdf2image pillow python-pptx
    Windows: también necesita poppler → https://github.com/oschwartz10612/poppler-windows/releases
    Descomprime poppler y pasa la ruta con --poppler "C:/poppler/Library/bin"
"""

import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import io

def pdf_to_pptx(pdf_path: str, output_path: str = None, dpi: int = 150, poppler_path: str = None):
    try:
        from pdf2image import convert_from_path
    except ImportError:
        print("ERROR: Instala pdf2image → pip install pdf2image")
        sys.exit(1)

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        print(f"ERROR: No se encontró el archivo: {pdf_path}")
        sys.exit(1)

    if output_path is None:
        output_path = pdf_path.with_suffix(".pptx")
    else:
        output_path = Path(output_path)

    print(f"Convirtiendo: {pdf_path.name}")
    print(f"DPI: {dpi}")

    # Convertir páginas PDF a imágenes
    kwargs = {"dpi": dpi}
    if poppler_path:
        kwargs["poppler_path"] = poppler_path

    try:
        pages = convert_from_path(str(pdf_path), **kwargs)
    except Exception as e:
        print(f"ERROR al leer el PDF: {e}")
        print("En Windows, asegúrate de tener poppler instalado y usa --poppler 'C:/poppler/Library/bin'")
        sys.exit(1)

    if not pages:
        print("ERROR: El PDF no tiene páginas.")
        sys.exit(1)

    print(f"Páginas encontradas: {len(pages)}")

    # Detectar tamaño de la primera página
    first = pages[0]
    w_px, h_px = first.size

    # Relación px → EMU (914400 EMU = 1 pulgada, a 'dpi' px/pulgada)
    px_to_emu = 914400 / dpi
    slide_w = int(w_px * px_to_emu)
    slide_h = int(h_px * px_to_emu)

    # Crear presentación con el tamaño exacto del PDF
    prs = Presentation()
    prs.slide_width = Emu(slide_w)
    prs.slide_height = Emu(slide_h)

    blank_layout = prs.slide_layouts[6]  # layout en blanco

    for i, page in enumerate(pages):
        slide = prs.slides.add_slide(blank_layout)

        # Guardar imagen en buffer
        img_buffer = io.BytesIO()
        page.save(img_buffer, format="PNG")
        img_buffer.seek(0)

        # Insertar imagen ocupando 100% del slide
        slide.shapes.add_picture(
            img_buffer,
            left=0, top=0,
            width=Emu(slide_w),
            height=Emu(slide_h)
        )

        print(f"  Slide {i + 1}/{len(pages)} ✓")

    prs.save(str(output_path))
    print(f"\n✅ PPTX guardado en: {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convierte PDF → PPTX (cada página = un slide con imagen)")
    parser.add_argument("pdf", help="Ruta al archivo PDF")
    parser.add_argument("-o", "--output", help="Ruta del PPTX de salida (opcional)")
    parser.add_argument("--dpi", type=int, default=150, help="Resolución de imagen (default: 150, recomendado: 150-200)")
    parser.add_argument("--poppler", help="Ruta a poppler/bin (solo Windows)")
    args = parser.parse_args()

    pdf_to_pptx(args.pdf, args.output, args.dpi, args.poppler)
